"""
Build the Brief 02 presentation deck.

Single output: B02_NilyufarShodmonova_INDEX_presentation.pptx (9 slides).
All visuals are constructed in python-pptx. No external image dependencies.
Audience is non-technical; slide copy stays in plain English.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# ---------- Design tokens ----------

NAVY = RGBColor(0x0B, 0x1B, 0x2E)
IVORY = RGBColor(0xF7, 0xF5, 0xEE)
ACCENT = RGBColor(0x4A, 0x88, 0xE5)
MUTED = RGBColor(0x6B, 0x86, 0xA6)
INK = RGBColor(0x1A, 0x1F, 0x2E)
INK_SOFT = RGBColor(0x2C, 0x3A, 0x55)
INK_MUTED = RGBColor(0x4A, 0x55, 0x70)
INK_FAINT = RGBColor(0x88, 0x97, 0xB1)
LINE = RGBColor(0xDC, 0xE2, 0xEE)
IVORY_INK_MUTED = RGBColor(0x6F, 0x6B, 0x5E)  # muted ink on ivory backgrounds

# Band gradient (from framework §06)
BAND_BELOW_BG = RGBColor(0xF2, 0xDE, 0xDB)
BAND_BELOW_FG = RGBColor(0xA2, 0x37, 0x10)
BAND_LOW_BG = RGBColor(0xFF, 0xEA, 0xD3)
BAND_LOW_FG = RGBColor(0x76, 0x61, 0x0C)
BAND_MED_BG = RGBColor(0xFF, 0xF8, 0xDA)
BAND_MED_FG = RGBColor(0x76, 0x61, 0x0C)
BAND_HIGH_BG = RGBColor(0xE2, 0xEE, 0xD9)
BAND_HIGH_FG = RGBColor(0x3F, 0x5A, 0x28)

FONT = "Helvetica"  # Falls back to Arial on Windows; both are clean sans.

# 16:9 widescreen.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- Text helpers ----------


def _set_letter_spacing(run, hundredths_of_pt: int) -> None:
    """Set letter spacing on a run via raw OOXML. python-pptx has no
    high-level wrapper for the `spc` attribute. Used for the small-caps look."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(hundredths_of_pt))


def _i(v) -> int:
    """Coerce an Emu / Inches / int / float to integer EMU. python-pptx's
    add_connector path writes raw floats into the OOXML attributes, which
    PowerPoint tolerates but Keynote and Google Slides reject. Anything
    crossing a connector boundary must be an integer."""
    return int(round(float(v)))


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=FONT,
    size=18,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
    spc=0,
    word_wrap=True,
):
    """Add a single-paragraph textbox. Returns the textbox shape."""
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if spc:
        _set_letter_spacing(run, spc)
    return tb


def add_lines(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    font=FONT,
    size=18,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.4,
    spc=0,
):
    """Multi-line textbox. Each entry of `lines` becomes its own paragraph."""
    tb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if spc:
            _set_letter_spacing(run, spc)
    return tb


# ---------- Shape helpers ----------


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill=None,
    line=None,
    line_weight=None,
    rounded=False,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        _i(x),
        _i(y),
        _i(w),
        _i(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_weight is not None:
            shape.line.width = Pt(line_weight)
    if rounded:
        # Tight corner radius. Overrides the default chunky rounded corners.
        shape.adjustments[0] = 0.10
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED, weight=0.75):
    line = slide.shapes.add_connector(1, _i(x1), _i(y1), _i(x2), _i(y2))  # 1 == straight
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_arrow(slide, x, y, w, h, color=ACCENT, weight=1.2):
    """Right-pointing arrow as a connector with arrow head."""
    cy = _i(y + h / 2)
    line = slide.shapes.add_connector(1, _i(x), cy, _i(x + w), cy)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # End arrow head via raw XML
    ln = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    from lxml import etree

    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def _set_dashed(shape) -> None:
    """Apply a dashed outline preset to a rect/line shape."""
    from pptx.oxml.ns import qn
    from lxml import etree

    ln = shape.line._get_or_add_ln()
    prstDash = ln.find(qn("a:prstDash"))
    if prstDash is None:
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")


# ---------- Common slide chrome ----------


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_wordmark(slide, on_navy: bool = False):
    """glomotec wordmark, top-left."""
    color = IVORY if on_navy else INK
    add_text(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(3.0),
        Inches(0.4),
        "glomotec",
        size=14,
        color=color,
        bold=False,
    )


def add_section_label(slide, x, y, w, text, *, on_navy: bool = False):
    color = MUTED
    add_text(
        slide,
        x,
        y,
        w,
        Inches(0.35),
        text.upper(),
        size=11,
        color=color,
        bold=True,
        spc=200,  # ~2pt of letter spacing for the small-caps feel
    )


def add_speaker_notes(slide, notes: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes
    for p in notes_tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(13)


# ---------- Slide builders ----------


def slide_1_title(prs):
    """Slide 1. Title (navy). UNCHANGED from prior version."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    add_text(
        s,
        Inches(0.6),
        Inches(0.45),
        Inches(3.0),
        Inches(0.4),
        "glomotec",
        size=14,
        color=IVORY,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(2.7),
        Inches(12.13),
        Inches(0.4),
        "FRONTIER PROGRAMME  ·  COHORT ONE  ·  BRIEF 02",
        size=12,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        spc=300,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(3.1),
        Inches(12.13),
        Inches(2.0),
        "INDEX",
        size=120,
        color=IVORY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        spc=80,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(5.05),
        Inches(12.13),
        Inches(0.6),
        "A system that reads UK immigration rules so advisors do not have to.",
        size=20,
        color=IVORY,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(6.6),
        Inches(12.13),
        Inches(0.4),
        "Nilyufar Shodmonova  ·  5 May 2026",
        size=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_speaker_notes(
        s,
        "Hi everyone. I'm Nilyufar. Brief 02. I'll walk you through INDEX, "
        "a system that reads UK immigration rules so advisors do not have "
        "to. I designed it and built a working version. Five minutes.",
    )


def slide_2_problem(prs):
    """Slide 2. The problem (ivory). Sets up why INDEX exists."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(8), "01  ·  THE PROBLEM")

    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.13),
        Inches(1.4),
        "Advisors read 80-page government documents from memory.",
        size=34,
        color=NAVY,
        bold=True,
        line_spacing=1.15,
    )

    # Single oversized stat replaces the bullet list. The presenter speaks
    # the bullets; the slide is a visual anchor for the price.
    add_text(
        s,
        Inches(0.6),
        Inches(3.7),
        Inches(12.13),
        Inches(2.0),
        "£8,000–£25,000 per case",
        size=88,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.05,
        spc=-30,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(5.7),
        Inches(12.13),
        Inches(0.45),
        "TYPICAL ADVISOR FEE",
        size=12,
        color=INK_FAINT,
        bold=True,
        align=PP_ALIGN.CENTER,
        spc=400,
    )

    # Speaker notes carry the original bullets so the presenter still narrates
    # the points the slide no longer prints.
    add_speaker_notes(
        s,
        "The UK Home Office publishes guidance for caseworkers across "
        "roughly 18 active immigration routes. These documents update "
        "every 1 to 4 months, without warning. A typical advisor charges "
        "between £8,000 and £25,000 per case. Failed applications cost "
        "clients time, money, and trust. glomotec advises on UK "
        "immigration. The rules are written for Home Office caseworkers, "
        "published online, and updated every few months without any "
        "notice. Most firms read each document once and trust their "
        "memory. Almost nobody tracks what changes between versions. "
        "That is the gap I am filling.",
    )


def slide_3_what_i_built(prs):
    """Slide 3. What I built (ivory). Plain-English summary."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(8), "02  ·  WHAT I BUILT")

    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.13),
        Inches(1.4),
        "INDEX reads the rules so advisors do not have to.",
        size=36,
        color=NAVY,
        bold=True,
        line_spacing=1.15,
    )

    # Three-step icon row replaces the bullet list. Each step is a circular
    # accent disc with a tiny label beneath. Arrows between the discs make
    # the read → structure → check flow legible at a glance.
    step_y = Inches(3.85)
    disc_d = Inches(0.95)
    label_h = Inches(0.45)
    n = 3
    margin_x = Inches(2.4)
    available = SLIDE_W - 2 * margin_x
    # Three discs with two arrows between. Disc-to-disc gap balances arrow length.
    gap = (available - n * disc_d) / (n - 1)

    steps = ["READ", "STRUCTURE", "CHECK"]
    centers = []
    for i, label in enumerate(steps):
        cx = margin_x + i * (disc_d + gap)
        # Disc.
        disc = s.shapes.add_shape(
            MSO_SHAPE.OVAL, _i(cx), _i(step_y), _i(disc_d), _i(disc_d),
        )
        disc.fill.solid()
        disc.fill.fore_color.rgb = RGBColor(0xEE, 0xF3, 0xFB)
        disc.line.color.rgb = ACCENT
        disc.line.width = Pt(1.0)
        # Numeral inside the disc.
        add_text(
            s,
            cx,
            step_y,
            disc_d,
            disc_d,
            str(i + 1),
            size=28,
            color=ACCENT,
            bold=True,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Label beneath, with a connecting hairline.
        add_text(
            s,
            cx - Inches(0.5),
            step_y + disc_d + Inches(0.25),
            disc_d + Inches(1.0),
            label_h,
            label,
            size=14,
            color=INK_SOFT,
            bold=True,
            align=PP_ALIGN.CENTER,
            spc=200,
        )
        centers.append(cx + disc_d / 2)

    # Arrows between each pair of discs.
    for i in range(n - 1):
        ax = margin_x + i * (disc_d + gap) + disc_d
        add_arrow(
            s,
            ax + Inches(0.05),
            step_y + disc_d / 2,
            gap - Inches(0.1),
            Inches(0),
            color=ACCENT,
            weight=1.4,
        )

    # Speaker notes carry the original bullets so the presenter narrates them.
    add_speaker_notes(
        s,
        "INDEX reads the rules so advisors do not have to. Three steps. "
        "It fetches every government guidance page. It turns long "
        "paragraphs into clear, separate rules. It checks a client "
        "against those rules. An advisor sees where the case is strong, "
        "where it is weak, what evidence is missing, and when a rule "
        "changes which of their clients is affected.",
    )


def slide_4_pipeline(prs):
    """Slide 4. How it works (ivory). Five modules with plain-English captions."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(8), "03  ·  HOW IT WORKS")

    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.13),
        Inches(0.9),
        "Five steps. One pipeline.",
        size=38,
        color=NAVY,
        bold=True,
    )

    # Five module boxes in a row with plain-English captions beneath each.
    modules = [
        ("CRAWLER", "downloads the rules", "from gov.uk"),
        ("EXTRACTOR", "turns long paragraphs", "into rules"),
        ("CHANGEFEED", "detects what", "changed between versions"),
        ("SCORER", "checks a client", "against the rules"),
        ("EVALUATOR", "keeps the system", "itself honest"),
    ]
    n = len(modules)
    margin_x = Inches(0.6)
    total_w = SLIDE_W - 2 * margin_x
    box_w = Inches(1.95)
    arrow_w = (total_w - n * box_w) / (n - 1)
    row_y = Inches(3.45)
    row_h = Inches(2.6)

    box_xs = [margin_x + i * (box_w + arrow_w) for i in range(n)]

    # Arrows behind boxes (between modules 1-4, 2-3, 3-4, 4-5).
    for i in range(n - 1):
        ax = box_xs[i] + box_w
        # Arrows aligned with the module-name vertical band, not the caption.
        add_arrow(
            s, ax, row_y + Inches(0.55), arrow_w, Inches(0.0), color=ACCENT, weight=1.5
        )

    for i, (name, line1, line2) in enumerate(modules):
        bx = box_xs[i]
        add_rect(
            s, bx, row_y, box_w, Inches(1.1),
            fill=RGBColor(0xEE, 0xF3, 0xFB),
            line=ACCENT,
            line_weight=0.9,
            rounded=True,
        )
        add_text(
            s,
            bx,
            row_y + Inches(0.32),
            box_w,
            Inches(0.5),
            name,
            size=14,
            color=ACCENT,
            bold=True,
            align=PP_ALIGN.CENTER,
            spc=120,
        )
        # Caption sits below the module box.
        cap_y = row_y + Inches(1.25)
        add_lines(
            s,
            bx,
            cap_y,
            box_w,
            Inches(1.1),
            [line1, line2],
            size=12,
            color=INK_MUTED,
            italic=True,
            align=PP_ALIGN.CENTER,
            line_spacing=1.3,
        )

    add_speaker_notes(
        s,
        "The first three steps run on a schedule. CRAWLER downloads every "
        "guidance page from gov.uk. EXTRACTOR uses an AI model to turn "
        "each long paragraph into a clear rule. CHANGEFEED compares the "
        "new version against the last one and tells the operations team "
        "what actually changed. SCORER runs whenever an advisor wants to "
        "check a client. EVALUATOR retests the whole system every night "
        "to catch any drop in quality before it reaches the team.",
    )


def slide_5_decision(prs):
    """Slide 5. The biggest decision (ivory). Text vs structured."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(10), "04  ·  THE BIGGEST DECISION")

    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.13),
        Inches(0.9),
        "Store rules as data, not as text.",
        size=40,
        color=NAVY,
        bold=True,
    )

    # Two-column comparison: Text (left) vs Structured (right).
    box_y = Inches(3.4)
    box_h = Inches(2.7)
    margin_x = Inches(0.6)
    gap = Inches(0.4)
    col_w = (SLIDE_W - 2 * margin_x - gap) / 2

    # Left: Text
    left_x = margin_x
    left_box = add_rect(
        s, left_x, box_y, col_w, box_h, fill=IVORY, line=MUTED, line_weight=0.75, rounded=True
    )
    _set_dashed(left_box)
    add_text(
        s,
        left_x + Inches(0.4),
        box_y + Inches(0.25),
        col_w - Inches(0.8),
        Inches(0.35),
        "TEXT",
        size=11,
        color=MUTED,
        bold=True,
        spc=200,
    )
    add_text(
        s,
        left_x + Inches(0.4),
        box_y + Inches(0.85),
        col_w - Inches(0.8),
        box_h - Inches(1.1),
        "“An applicant must be 18 or over and must have endorsement from an approved body.”",
        size=18,
        color=INK_SOFT,
        italic=True,
        line_spacing=1.4,
    )

    # Right: Structured
    right_x = margin_x + col_w + gap
    add_rect(
        s, right_x, box_y, col_w, box_h,
        fill=RGBColor(0xEE, 0xF3, 0xFB),
        line=ACCENT,
        line_weight=1.0,
        rounded=True,
    )
    add_text(
        s,
        right_x + Inches(0.4),
        box_y + Inches(0.25),
        col_w - Inches(0.8),
        Inches(0.35),
        "STRUCTURED",
        size=11,
        color=ACCENT,
        bold=True,
        spc=200,
    )
    add_lines(
        s,
        right_x + Inches(0.4),
        box_y + Inches(0.8),
        col_w - Inches(0.8),
        box_h - Inches(1.0),
        [
            "Two separate rules. Each tagged with:",
            "·  who must satisfy it",
            "·  what evidence proves it",
            "·  what the threshold is",
            "·  what stage of the journey it applies at",
        ],
        size=15,
        color=INK_SOFT,
        line_spacing=1.45,
    )

    # Subline below the columns.
    add_text(
        s,
        Inches(0.6),
        Inches(6.55),
        Inches(12.13),
        Inches(0.55),
        "Text can be read. Only structure can be checked, compared across versions, and updated automatically.",
        size=14,
        color=INK_MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_speaker_notes(
        s,
        "I had three options. Use AI to read the rules and answer questions "
        "like a chatbot. Summarise the rules in plain prose. Or pull every "
        "rule out into a clear, separate record. The first two cannot "
        "reliably catch what changes between versions, and a single word "
        "added or removed in government guidance can invalidate hundreds "
        "of applications. Only the structured version catches that. So I "
        "chose structured. Each rule is stored on its own with what kind "
        "of test it is, who has to prove it, and what evidence counts. "
        "That is what lets the system know that \"must be 18 or over\" "
        "and \"must have a viable business\" need very different evidence "
        "even though both end in a yes or no.",
    )


def slide_6_scoring(prs):
    """Slide 6. How scoring works (ivory). Three readiness axes + four bands."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(8), "05  ·  HOW SCORING WORKS")

    add_text(
        s,
        Inches(0.6),
        Inches(1.85),
        Inches(12.13),
        Inches(0.85),
        "Three things we check. Four levels of confidence.",
        size=32,
        color=NAVY,
        bold=True,
        line_spacing=1.15,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(2.7),
        Inches(12.13),
        Inches(0.45),
        "We do not show numbers. We show recommendations.",
        size=16,
        color=INK_MUTED,
        line_spacing=1.2,
    )

    # ---------- Top half: three readiness axis cards ----------
    axis_y = Inches(3.25)
    axis_h = Inches(1.45)
    margin_x = Inches(0.6)
    gap = Inches(0.3)
    card_w = (SLIDE_W - 2 * margin_x - 2 * gap) / 3

    # Suitability accent: deeper amber, mirrors the verdict-hero column label.
    SUITABILITY_FG = RGBColor(0x9A, 0x45, 0x15)
    SUITABILITY_BG = RGBColor(0xFB, 0xEF, 0xE4)

    axes = [
        (
            "FIT",
            "Does the applicant fit the route?",
            ACCENT,
            RGBColor(0xEE, 0xF3, 0xFB),
        ),
        (
            "PAPERWORK",
            "Is the application complete?",
            ACCENT,
            RGBColor(0xEE, 0xF3, 0xFB),
        ),
        (
            "SUITABILITY",
            "Could the case be refused for another reason?",
            SUITABILITY_FG,
            SUITABILITY_BG,
        ),
    ]
    for i, (label, blurb, fg, bg) in enumerate(axes):
        cx = margin_x + i * (card_w + gap)
        add_rect(
            s, cx, axis_y, card_w, axis_h,
            fill=bg, line=fg, line_weight=1.0, rounded=True,
        )
        add_text(
            s,
            cx + Inches(0.3),
            axis_y + Inches(0.28),
            card_w - Inches(0.6),
            Inches(0.4),
            label,
            size=12,
            color=fg,
            bold=True,
            spc=200,
        )
        add_text(
            s,
            cx + Inches(0.3),
            axis_y + Inches(0.78),
            card_w - Inches(0.6),
            Inches(0.6),
            blurb,
            size=15,
            color=INK_SOFT,
            line_spacing=1.35,
        )

    # ---------- Bottom half: four-band gradient ----------
    band_y = Inches(5.2)
    band_h = Inches(0.7)
    band_margin_x = Inches(0.6)
    band_w = (SLIDE_W - 2 * band_margin_x) / 4
    bands = [
        ("BELOW THRESHOLD", "Do not submit", BAND_BELOW_BG, BAND_BELOW_FG),
        ("LOW", "Get more evidence", BAND_LOW_BG, BAND_LOW_FG),
        ("MEDIUM", "Advisor review", BAND_MED_BG, BAND_MED_FG),
        ("HIGH", "Recommend submission", BAND_HIGH_BG, BAND_HIGH_FG),
    ]
    for i, (label, action, bg, fg) in enumerate(bands):
        bx = band_margin_x + i * band_w
        add_rect(s, bx, band_y, band_w, band_h, fill=bg, line=fg, line_weight=0.5)
        add_text(
            s,
            bx,
            band_y + Inches(0.13),
            band_w,
            Inches(0.4),
            label,
            size=11,
            color=fg,
            bold=True,
            align=PP_ALIGN.CENTER,
            spc=200,
        )
        add_text(
            s,
            bx,
            band_y + band_h + Inches(0.12),
            band_w,
            Inches(0.4),
            action,
            size=13,
            color=INK_SOFT,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    add_speaker_notes(
        s,
        "INDEX checks a client on three things. Fit answers does the "
        "applicant actually fit this visa route. Paperwork answers is the "
        "application complete and properly evidenced. Suitability answers "
        "could this case still be refused for another reason, things like "
        "past immigration issues, a criminal record, false statements in a "
        "previous application, or unpaid NHS debt. The result on each "
        "check sits in one of four levels of confidence, and each level "
        "comes with a clear action. Advisors see the level and the action, "
        "not a raw number.",
    )


def slide_7_demo(prs):
    """Slide 7. Live demo (navy). UNCHANGED."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    add_text(
        s,
        Inches(0.6),
        Inches(0.45),
        Inches(3.0),
        Inches(0.4),
        "glomotec",
        size=14,
        color=IVORY,
    )

    add_section_label(
        s, Inches(0.6), Inches(2.5), Inches(12.13), "LIVE", on_navy=True
    )

    add_text(
        s,
        Inches(0.6),
        Inches(3.05),
        Inches(12.13),
        Inches(1.6),
        "index-advisor.vercel.app",
        size=64,
        color=IVORY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Two-line subline. Line 1 names the framing; line 2 names the surfaces.
    add_text(
        s,
        Inches(0.6),
        Inches(4.85),
        Inches(12.13),
        Inches(0.5),
        "One system, two surfaces.",
        size=22,
        color=IVORY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        Inches(0.6),
        Inches(5.4),
        Inches(12.13),
        Inches(0.5),
        "For advisors. For new prospects.",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    # Footer: two side-by-side groupings, separated by ample whitespace.
    # Left = ADVISOR · /clients (full pipeline). Right = PROSPECT · /signal
    # (substantive-only SCORER). The whitespace between them carries the
    # "two surfaces, one engine" framing visually.
    footer_y = Inches(6.5)
    margin_x = Inches(0.6)
    gap = Inches(1.6)
    col_w = (SLIDE_W - 2 * margin_x - gap) / 2

    # Left: ADVISOR · /clients
    left_x = margin_x
    add_text(
        s,
        left_x,
        footer_y,
        col_w,
        Inches(0.35),
        "FOR ADVISORS  ·  /clients",
        size=11,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
        spc=300,
    )
    add_text(
        s,
        left_x,
        footer_y + Inches(0.4),
        col_w,
        Inches(0.35),
        "CRAWLER  ·  EXTRACTOR  ·  CHANGEFEED  ·  SCORER  ·  EVALUATOR",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        spc=200,
    )

    # Right: PROSPECT · /signal
    right_x = margin_x + col_w + gap
    add_text(
        s,
        right_x,
        footer_y,
        col_w,
        Inches(0.35),
        "FOR PROSPECTS  ·  /signal",
        size=11,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
        spc=300,
    )
    add_text(
        s,
        right_x,
        footer_y + Inches(0.4),
        col_w,
        Inches(0.35),
        "SCORER (substantive only)",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        spc=200,
    )

    add_speaker_notes(
        s,
        "I will switch to the browser. First I will click Run live "
        "pipeline. CRAWLER downloads from gov.uk live, CHANGEFEED spots "
        "what changed against the last version, EXTRACTOR pulls one rule "
        "through the AI model. About fifteen seconds. Then I will click "
        "into a saved client to show the three checks, including the new "
        "suitability one. Then over to /signal, the version for new "
        "prospects, a quick chat as a Brazilian fintech founder. The "
        "conversation ends with a Strong fit verdict. Same system, two "
        "surfaces. About ninety seconds total.",
    )


def slide_8_surprise(prs):
    """Slide 8. What I learned (ivory). The operator-consultation story."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    add_wordmark(s)

    add_section_label(s, Inches(0.6), Inches(1.4), Inches(8), "06  ·  WHAT I LEARNED")

    add_text(
        s,
        Inches(0.6),
        Inches(1.9),
        Inches(12.13),
        Inches(1.05),
        "I asked the person who runs this work every day.",
        size=30,
        color=NAVY,
        bold=True,
        line_spacing=1.15,
    )

    # Single body line replaces the bullet list. Centred, large, the slide
    # is now a visual anchor; the presenter speaks the rest.
    add_text(
        s,
        Inches(0.6),
        Inches(3.6),
        Inches(12.13),
        Inches(1.4),
        "She said: add the third check.",
        size=40,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.2,
    )

    # Three-pillar diagram below the body. Suitability uses the verdict-hero
    # accent (deeper amber) to signal that it is the gate axis.
    SUITABILITY_FG = RGBColor(0x9A, 0x45, 0x15)
    SUITABILITY_BG = RGBColor(0xFB, 0xEF, 0xE4)

    pillar_y = Inches(6.05)
    pillar_h = Inches(0.7)
    pillar_margin_x = Inches(2.0)
    pillar_gap = Inches(0.4)
    pillar_w = (SLIDE_W - 2 * pillar_margin_x - 2 * pillar_gap) / 3
    pillars = [
        ("FIT", ACCENT, RGBColor(0xEE, 0xF3, 0xFB)),
        ("PAPERWORK", ACCENT, RGBColor(0xEE, 0xF3, 0xFB)),
        ("SUITABILITY", SUITABILITY_FG, SUITABILITY_BG),
    ]
    for i, (label, fg, bg) in enumerate(pillars):
        px = pillar_margin_x + i * (pillar_w + pillar_gap)
        add_rect(
            s, px, pillar_y, pillar_w, pillar_h,
            fill=bg, line=fg, line_weight=1.0, rounded=True,
        )
        add_text(
            s,
            px,
            pillar_y + Inches(0.16),
            pillar_w,
            Inches(0.45),
            label,
            size=15,
            color=fg,
            bold=True,
            align=PP_ALIGN.CENTER,
            spc=160,
        )

    add_speaker_notes(
        s,
        "The most useful thing I did this week was not more testing. It "
        "was sending my work to Nameera, the Director of Operations at "
        "glomotec. She runs the team that actually files these "
        "applications, so she sees every way a case can fall apart in "
        "real life. She validated the first two checks, fit and "
        "paperwork. Then she told me a case can still be refused even "
        "when the person fits the route and the paperwork is complete. "
        "For example, past immigration issues, a criminal record, false "
        "statements in a previous application, or unpaid NHS debt. That "
        "is a third kind of check, and I had been treating only the "
        "first two. So I added it. Three checks now, not two, and the "
        "document, the code, and the live demo all reflect that. Ray, "
        "the CEO, also wrote back. He pointed out the same system could "
        "help new people who are not yet clients, not just the ones we "
        "already have. Same rules, different lens. So I built that "
        "surface too, the prospect-facing version called /signal. You "
        "will see it in the demo. The things that mattered most this "
        "week came from talking to the people who actually do this "
        "work.",
    )


def slide_9_thanks(prs):
    """Slide 9. Questions (navy). UNCHANGED."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)

    add_text(
        s,
        Inches(0.6),
        Inches(0.45),
        Inches(3.0),
        Inches(0.4),
        "glomotec",
        size=14,
        color=IVORY,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(2.65),
        Inches(12.13),
        Inches(0.4),
        "THANK YOU",
        size=14,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        spc=400,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(3.2),
        Inches(12.13),
        Inches(2.0),
        "Questions?",
        size=88,
        color=IVORY,
        bold=True,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_text(
        s,
        Inches(0.6),
        Inches(6.55),
        Inches(12.13),
        Inches(0.4),
        "Nilyufar Shodmonova  ·  Brief 02  ·  INDEX  ·  github.com/lotirium/glomotec-index",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_speaker_notes(
        s,
        "That's INDEX. The framework PDF and the code are in the repo. "
        "Happy to take questions.",
    )


# ---------- Main ----------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_what_i_built(prs)
    slide_4_pipeline(prs)
    slide_5_decision(prs)
    slide_6_scoring(prs)
    slide_7_demo(prs)
    slide_8_surprise(prs)
    slide_9_thanks(prs)

    out = "B02_NilyufarShodmonova_INDEX_presentation.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    out = build()
    from pptx import Presentation as _P

    n = len(_P(out).slides)
    print(f"OK  {out}  slides={n}")
